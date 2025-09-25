import os
import pdfplumber
import docx
import pptx
import markdown
from pathlib import Path
from typing import List, Dict, Tuple
from config import config
from utils.logging import logger

class DocumentLoader:
    """文档加载器，支持多种格式的文档加载和内容提取"""
    
    def __init__(self):
        """初始化文档加载器"""
        self.supported_extensions = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".doc": self._load_docx,  # 简化处理，实际可能需要不同方法
            ".pptx": self._load_pptx,
            ".md": self._load_markdown,
            ".txt": self._load_text
        }
        logger.info(f"文档加载器初始化完成，支持格式: {list(self.supported_extensions.keys())}")
    
    def load_document(self, file_path: str or Path) -> Tuple[str, Dict]:
        """
        加载单个文档并提取内容和元数据
        
        Args:
            file_path: 文档路径
            
        Returns:
            文档内容和元数据的元组
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
            
        # 获取文件扩展名
        ext = file_path.suffix.lower()
        
        if ext not in self.supported_extensions:
            raise ValueError(f"不支持的文件格式: {ext}，支持的格式: {list(self.supported_extensions.keys())}")
        
        # 提取元数据
        metadata = {
            "file_name": file_path.name,
            "file_path": str(file_path),
            "file_size": os.path.getsize(file_path),
            "file_type": ext[1:],  # 去除点号
            "last_modified": os.path.getmtime(file_path)
        }
        
        # 调用对应格式的加载方法
        logger.info(f"加载文档: {file_path.name}，格式: {ext}")
        content = self.supported_extensions[ext](file_path)
        
        return content, metadata
    
    def load_directory(self, dir_path: str or Path) -> List[Tuple[str, Dict]]:
        """
        加载目录中的所有支持的文档
        
        Args:
            dir_path: 目录路径
            
        Returns:
            文档内容和元数据的列表
        """
        dir_path = Path(dir_path)
        if not dir_path.is_dir():
            raise NotADirectoryError(f"不是有效的目录: {dir_path}")
            
        documents = []
        for file_path in dir_path.rglob("*"):
            if file_path.is_file():
                ext = file_path.suffix.lower()
                if ext in self.supported_extensions:
                    try:
                        content, metadata = self.load_document(file_path)
                        documents.append((content, metadata))
                    except Exception as e:
                        logger.error(f"加载文档 {file_path} 失败: {str(e)}", exc_info=True)
        
        logger.info(f"从目录 {dir_path} 加载了 {len(documents)} 个文档")
        return documents
    
    def _load_pdf(self, file_path: Path) -> str:
        """加载PDF文件内容"""
        content = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    content.append(f"Page {page_num}:\n{text}")
        
        return "\n\n".join(content)
    
    def _load_docx(self, file_path: Path) -> str:
        """加载DOCX文件内容"""
        doc = docx.Document(file_path)
        content = []
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
        
        return "\n\n".join(content)
    
    def _load_pptx(self, file_path: Path) -> str:
        """加载PPTX文件内容"""
        prs = pptx.Presentation(file_path)
        content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            if slide_content:
                content.append(f"Slide {slide_num}:\n" + "\n\n".join(slide_content))
        
        return "\n\n".join(content)
    
    def _load_markdown(self, file_path: Path) -> str:
        """加载Markdown文件内容（转为纯文本）"""
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()
        
        # 转换为纯文本（去除Markdown格式）
        html = markdown.markdown(md_content)
        # 简单处理HTML标签，实际应用可能需要更复杂的处理
        import re
        text = re.sub(r"<.*?>", "", html)
        return text
    
    def _load_text(self, file_path: Path) -> str:
        """加载文本文件内容"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()

# 创建文档加载器实例
document_loader = DocumentLoader()
